"""
Build the progress-update slide deck for the 27 May 2026 meeting
from the matching meeting_prep notes.

Outputs: meeting_prep/2026-05-27_progress_slides.pptx
Image source: meeting_prep/slide_figures/  (compressed overlays)

Style:
- Warm Terracotta palette (paper-themed)
- Georgia headers + Calibri body
- Big-number callouts where the result is a number
- Real overlay images where the result is visual
- Mix of layouts so no two consecutive slides look the same
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "meeting_prep" / "2026-05-27_progress_slides.pptx"
FIG = ROOT / "meeting_prep" / "slide_figures"
OUT.parent.mkdir(parents=True, exist_ok=True)

# --- Palette --------------------------------------------------------------
TERRACOTTA = RGBColor(0xB8, 0x50, 0x42)
SAND       = RGBColor(0xE7, 0xE8, 0xD1)
SAGE       = RGBColor(0xA7, 0xBE, 0xAE)
INK        = RGBColor(0x2A, 0x1F, 0x1A)
INK_LIGHT  = RGBColor(0x6B, 0x5A, 0x52)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

HEADER_FONT = "Georgia"
BODY_FONT   = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)


def add_bg(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_text(slide, left, top, width, height, text, *,
             font=BODY_FONT, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=15, color=INK, line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•   {item}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
    rect.shadow.inherit = False
    return rect


def add_round_rect(slide, left, top, width, height, fill, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.08
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_image(slide, path, left, top, width=None, height=None):
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


# ==========================================================================
# Slide 1: title
# ==========================================================================
s = add_slide()
add_bg(s, TERRACOTTA)
add_text(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.4),
         "PROGRESS UPDATE  /  27 MAY 2026",
         font=HEADER_FONT, size=14, color=SAND)
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6),
         "Quantitative laid-line analysis\nfrom DRP and MSI imagery",
         font=HEADER_FONT, size=44, color=SAND, line_spacing=1.15)
add_text(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.8),
         "A progress report for Prof Orietta Da Rold",
         font=HEADER_FONT, size=20, color=SAND)
add_text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.2),
         "Zhanfeng Zhou\nIIB project, Department of Engineering\nSupervisor: TODO",
         font=BODY_FONT, size=14, color=SAND, line_spacing=1.3)
add_rect(s, Inches(12.5), Inches(0), Inches(0.83), SLIDE_H, SAND)
add_text(s, Inches(12.55), Inches(0.4), Inches(0.7), Inches(0.4),
         "01", font=HEADER_FONT, size=22, color=TERRACOTTA, bold=True)


# ==========================================================================
# Slide 2: where we are
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
         "Where we are",
         font=HEADER_FONT, size=36, color=INK)
add_text(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.5),
         "The problem", font=HEADER_FONT, size=20, color=TERRACOTTA)
add_text(s, Inches(0.7), Inches(2.05), Inches(5.8), Inches(4.5),
         "Codicology uses laid-line density and chain-line spacing to "
         "identify the moulds on which medieval and early-modern paper "
         "was formed. The standard measurement is by hand under raking "
         "light: slow, observer-dependent, and tiring on fragile material. "
         "Multispectral and DRM scans now exist at scale, but until now "
         "there has been no automated way to turn those scans into the "
         "quantitative density and wire-diameter figures that "
         "mould-pair identification needs.",
         font=BODY_FONT, size=14, color=INK, line_spacing=1.35)
add_text(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.5),
         "What's been built", font=HEADER_FONT, size=20, color=TERRACOTTA)
add_bullets(s, Inches(7.0), Inches(2.05), Inches(5.6), Inches(4.5),
            [
                "A two-variant detection pipeline on a shared framework",
                "Multi-phi variant for full DRP stacks",
                "Single-image variant for legacy MSI scans",
                "Per-folio metric bundle for confidence flagging",
                "Validated on synthetic phantoms and on real folios",
            ], size=15)
add_rect(s, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), TERRACOTTA)
add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.35),
         "02  /  progress update for Prof Da Rold  /  27 May 2026",
         font=BODY_FONT, size=10, color=SAND)


# ==========================================================================
# Slide 3: pipeline overview
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
         "Two inputs, two detectors, one framework",
         font=HEADER_FONT, size=32, color=INK)
box_w = Inches(4.2); box_h = Inches(0.9)
in_y = Inches(1.6)
add_round_rect(s, Inches(1.0), in_y, box_w, box_h, SAGE)
add_text(s, Inches(1.0), in_y, box_w, box_h, "DRP stack (N images)",
         font=BODY_FONT, size=16, color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_round_rect(s, Inches(8.1), in_y, box_w, box_h, SAGE)
add_text(s, Inches(8.1), in_y, box_w, box_h, "Single MSI image",
         font=BODY_FONT, size=16, color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pp_y = Inches(3.0)
add_round_rect(s, Inches(1.0), pp_y, Inches(11.3), Inches(1.1), TERRACOTTA)
add_text(s, Inches(1.0), pp_y, Inches(11.3), Inches(1.1),
         "Common preprocessing\nbackground subtraction  /  auto line direction  /  polarity flag",
         font=BODY_FONT, size=14, color=SAND, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
det_y = Inches(4.7)
add_round_rect(s, Inches(1.0), det_y, box_w, Inches(1.2), WHITE, line=TERRACOTTA)
add_text(s, Inches(1.0), det_y, box_w, Inches(1.2),
         "Multi-phi DRP variant\naggregated spectra  +  phase consensus",
         font=BODY_FONT, size=13, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
add_round_rect(s, Inches(8.1), det_y, box_w, Inches(1.2), WHITE, line=TERRACOTTA)
add_text(s, Inches(8.1), det_y, box_w, Inches(1.2),
         "Single-image MSI variant\nradial FFT  +  Gabor cleanup",
         font=BODY_FONT, size=13, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
add_round_rect(s, Inches(2.6), Inches(6.3), Inches(8.0), Inches(0.7), INK)
add_text(s, Inches(2.6), Inches(6.3), Inches(8.0), Inches(0.7),
         "Density  /  wire FWHM  /  metric bundle (4 or 5)",
         font=BODY_FONT, size=14, color=SAND, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================================================
# Slide 4: calibrated MSI headline + f5v overlay
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.7),
         "Headline result on the calibrated pair",
         font=HEADER_FONT, size=30, color=INK)
add_text(s, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.5),
         "Kk.1.5 f5v and f9v: MSI transmitted-light, your manual count = 9 lines/cm",
         font=BODY_FONT, size=14, color=INK_LIGHT)

# Left: f5v overlay image
img_x = Inches(0.5)
img_y = Inches(1.7)
img_h = Inches(5.0)
add_image(s, FIG / "msi_f5v_good.jpeg", img_x, img_y, height=img_h)
add_text(s, img_x, Inches(6.75), Inches(5.0), Inches(0.4),
         "Pipeline output on Kk.1.5 f5v (red grid = detected wires)",
         font=BODY_FONT, size=10, color=INK_LIGHT,
         align=PP_ALIGN.LEFT)

# Right: two stacked callouts
def small_callout(slide, left, top, width, height, label, pipeline, err):
    add_round_rect(slide, left, top, width, height, WHITE, line=SAGE)
    pad = Inches(0.2)
    add_text(slide, left + pad, top + Inches(0.1), width - pad*2, Inches(0.35),
             label, font=HEADER_FONT, size=14, color=TERRACOTTA, bold=True)
    add_text(slide, left + pad, top + Inches(0.5), Inches(2.2), Inches(0.9),
             pipeline, font=HEADER_FONT, size=44, color=INK, bold=True)
    add_text(slide, left + pad + Inches(2.4), top + Inches(0.62),
             width - pad*2 - Inches(2.4), Inches(0.4),
             "lines / cm", font=BODY_FONT, size=12, color=INK_LIGHT)
    add_text(slide, left + pad + Inches(2.4), top + Inches(0.95),
             width - pad*2 - Inches(2.4), Inches(0.4),
             f"manual 9.0  |  error {err}",
             font=BODY_FONT, size=12, color=INK)

right_x = Inches(6.5)
right_w = Inches(6.4)
small_callout(s, right_x, Inches(1.7), right_w, Inches(1.5),
              "Kk.1.5  f5v", "8.97", "−0.35 %")
small_callout(s, right_x, Inches(3.35), right_w, Inches(1.5),
              "Kk.1.5  f9v", "9.04", "+0.48 %")

# Caveat below the callouts
add_text(s, right_x, Inches(5.05), right_w, Inches(2.0),
         "Within the intrinsic ±1 wire/cm precision of a manual count "
         "over a 1 cm window. What this shows: pipeline introduces no "
         "bias the manual reference can detect. What it does not show: "
         "absolute accuracy below the manual reference's precision floor.",
         font=BODY_FONT, size=12, color=INK_LIGHT, line_spacing=1.35)


# ==========================================================================
# Slide 5: pipeline on other MSI folios (success + failure side by side)
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.7),
         "What the output looks like on other MSI folios",
         font=HEADER_FONT, size=28, color=INK)

# Left half: a success case (Hh2-12)
img_y = Inches(1.5)
img_h = Inches(4.6)
half_w = Inches(6.0)
left_x = Inches(0.5)
right_x = Inches(6.8)

add_image(s, FIG / "msi_hh2_decent.jpeg", left_x + Inches(1.2), img_y, height=img_h)
add_text(s, left_x, Inches(6.25), half_w, Inches(0.4),
         "Hh.2.12 f190 (working case)",
         font=HEADER_FONT, size=15, color=INK, bold=True)
add_text(s, left_x, Inches(6.65), half_w, Inches(0.8),
         "Pipeline: 8.91 lines/cm  /  spreadsheet 10.0  /  self-z = +5.6",
         font=BODY_FONT, size=12, color=INK_LIGHT)
add_text(s, left_x, Inches(7.0), half_w, Inches(0.4),
         "Density consistent; the 10% spreadsheet gap is within-stock variation.",
         font=BODY_FONT, size=11, color=INK_LIGHT)

# Right half: failure case (Ff4-15)
add_image(s, FIG / "msi_ff4_15_fail.jpeg", right_x + Inches(1.2), img_y, height=img_h)
add_text(s, right_x, Inches(6.25), half_w, Inches(0.4),
         "Ff.4.15 f24r (sub-harmonic alias)",
         font=HEADER_FONT, size=15, color=INK, bold=True)
add_text(s, right_x, Inches(6.65), half_w, Inches(0.8),
         "Pipeline: 3.92 lines/cm  /  spreadsheet 13.5  /  self-z = −0.8",
         font=BODY_FONT, size=12, color=INK_LIGHT)
add_text(s, right_x, Inches(7.0), half_w, Inches(0.4),
         "Grid clearly too sparse; the metric bundle flags it: do not trust this density.",
         font=BODY_FONT, size=11, color=TERRACOTTA, bold=True)


# ==========================================================================
# Slide 6: spreadsheet discrepancy (conceptual)
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
         "The spreadsheet discrepancy",
         font=HEADER_FONT, size=30, color=INK)
add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         "Pipeline disagrees with the catalogue on 7 of 9 folios, by 10 – 71 per cent",
         font=BODY_FONT, size=14, color=INK_LIGHT)

def column(slide, left, top, width, header, header_color, body):
    add_text(slide, left, top, width, Inches(0.5),
             header, font=HEADER_FONT, size=18, color=header_color, bold=True)
    add_text(slide, left, top + Inches(0.6), width, Inches(4.5),
             body, font=BODY_FONT, size=13, color=INK, line_spacing=1.4)

column(s, Inches(0.7), Inches(2.0), Inches(5.8),
       "Face-value reading", TERRACOTTA,
       "A systematic detector under-count would be alarming, and would make the "
       "pipeline useless as a forensic identifier.\n\n"
       "But this reading is ruled out by the two manual-GT folios: pipeline and "
       "manual count agree with each other and disagree with the spreadsheet by "
       "the same amount.")

column(s, Inches(7.0), Inches(2.0), Inches(5.7),
       "Working hypothesis", SAGE,
       "The spreadsheet records one density per stock, measured on one folio. "
       "Within a stock, different folios can come from different mould pairs in "
       "the same papermaker batch.\n\n"
       "If correct, the seven errors are a measurement of within-stock variation, "
       "not detector noise.")

add_round_rect(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.85), INK)
add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.85),
         "Does the within-stock hypothesis match what you see in the corpus?",
         font=HEADER_FONT, size=18, color=SAND, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================================================
# Slide 7: dataset 10 (Da Rold reproduction)
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.7),
         "Your hand-made reproduction sample",
         font=HEADER_FONT, size=30, color=INK)
add_text(s, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.5),
         "Dataset 10: paper made by hand from a reproduction of a 15th-century European mould",
         font=BODY_FONT, size=13, color=INK_LIGHT)

# Image (landscape, fills upper area)
add_image(s, FIG / "drp_dataset10_darold.jpeg",
          Inches(0.5), Inches(1.7), width=Inches(7.5))
add_text(s, Inches(0.5), Inches(5.85), Inches(7.5), Inches(0.4),
         "Pipeline output on dataset 10 (red grid = detected wires)",
         font=BODY_FONT, size=10, color=INK_LIGHT)

# Right column: number + metric bundle
rx = Inches(8.4)
rw = Inches(4.5)
add_round_rect(s, rx, Inches(1.7), rw, Inches(2.4), WHITE, line=SAGE)
add_text(s, rx + Inches(0.2), Inches(1.85), rw - Inches(0.4), Inches(0.35),
         "PIPELINE OUTPUT",
         font=HEADER_FONT, size=12, color=TERRACOTTA, bold=True)
add_text(s, rx + Inches(0.2), Inches(2.25), rw - Inches(0.4), Inches(1.3),
         "8.31", font=HEADER_FONT, size=64, color=INK, bold=True)
add_text(s, rx + Inches(0.2), Inches(3.55), rw - Inches(0.4), Inches(0.3),
         "lines per centimetre",
         font=BODY_FONT, size=12, color=INK_LIGHT)
add_text(s, rx + Inches(0.2), Inches(3.8), rw - Inches(0.4), Inches(0.3),
         "wire FWHM 0.38 mm at the median",
         font=BODY_FONT, size=11, color=INK)

# Metric bundle below the big callout
add_text(s, rx, Inches(4.4), rw, Inches(0.4),
         "Metric bundle: all five healthy",
         font=HEADER_FONT, size=14, color=TERRACOTTA, bold=True)

def metric_row(slide, left, top, label, value):
    add_text(slide, left, top, Inches(3.0), Inches(0.32),
             label, font=BODY_FONT, size=11, color=INK)
    add_text(slide, left + Inches(3.0), top, Inches(1.5), Inches(0.32),
             value, font=BODY_FONT, size=11, color=INK, bold=True)

metric_row(s, rx, Inches(4.85), "Self-contrast z",     "+4.05")
metric_row(s, rx, Inches(5.15), "Harmonic-fit R²",     "0.35")
metric_row(s, rx, Inches(5.45), "Split-half period",   "< 0.001 px")
metric_row(s, rx, Inches(5.75), "Wire FWHM IQR",       "0.33 – 0.43 mm")
metric_row(s, rx, Inches(6.05), "Interval CV",         "14 %")

# Ask box
add_round_rect(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.75), INK)
add_text(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.75),
         "Could you share the laid density of the mould? It turns this into a true ground truth.",
         font=HEADER_FONT, size=15, color=SAND, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================================================
# Slide 8: DRP zoomed-in on real folios + crop/band recipe
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.7),
         "DRP on zoomed-in real folios",
         font=HEADER_FONT, size=28, color=INK)
add_text(s, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.5),
         "Acquired at ~8.65 cm field of view; covers a small region of a folio at higher resolution",
         font=BODY_FONT, size=13, color=INK_LIGHT)

# Top: dataset 4 (success case)
add_image(s, FIG / "drp_dataset4_zoomed.jpeg",
          Inches(0.5), Inches(1.7), width=Inches(7.0))
add_text(s, Inches(7.7), Inches(1.7), Inches(5.2), Inches(0.4),
         "Dataset 4 (success)",
         font=HEADER_FONT, size=15, color=INK, bold=True)
add_text(s, Inches(7.7), Inches(2.1), Inches(5.2), Inches(2.0),
         "Pipeline: 11.3 lines/cm  /  self-z = +2.1\n\n"
         "Narrowing the period search range to 8 – 50 px excludes "
         "a long-period alias and gives a clean detection. Wire "
         "diameter and density both consistent with handmade paper "
         "from the period.",
         font=BODY_FONT, size=12, color=INK, line_spacing=1.4)

# Bottom: dataset 2 vs 2b — same source, different crop
add_image(s, FIG / "drp_dataset2_fullframe.jpeg",
          Inches(0.5), Inches(4.55), width=Inches(5.5))
add_text(s, Inches(0.5), Inches(7.05), Inches(5.5), Inches(0.4),
         "Dataset 2 (full frame)  —  self-z = +0.36 (low confidence)",
         font=BODY_FONT, size=11, color=INK_LIGHT)

add_image(s, FIG / "drp_dataset2b_cropped.jpeg",
          Inches(6.3), Inches(4.55), width=Inches(5.5))
add_text(s, Inches(6.3), Inches(7.05), Inches(5.5), Inches(0.4),
         "Dataset 2b (cropped to clean region)  —  self-z = +3.35",
         font=BODY_FONT, size=11, color=TERRACOTTA, bold=True)

# Right-side caption for the comparison
add_text(s, Inches(11.85), Inches(4.55), Inches(1.4), Inches(2.5),
         "Same source\nbut a small crop\nlifts confidence\nten-fold.",
         font=HEADER_FONT, size=13, color=INK, line_spacing=1.4)


# ==========================================================================
# Slide 9: three asks
# ==========================================================================
s = add_slide()
add_bg(s, SAND)
add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
         "What I would value from today's meeting",
         font=HEADER_FONT, size=32, color=INK)

def ask_card(slide, left, top, width, height, num, title, body):
    add_round_rect(slide, left, top, width, height, WHITE, line=SAGE)
    badge_d = Inches(0.95)
    add_round_rect(slide, left + Inches(0.3), top + Inches(0.3), badge_d, badge_d,
                   TERRACOTTA)
    add_text(slide, left + Inches(0.3), top + Inches(0.3), badge_d, badge_d,
             num, font=HEADER_FONT, size=32, color=SAND, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_x = left + Inches(1.5)
    text_w = width - Inches(1.8)
    add_text(slide, text_x, top + Inches(0.35), text_w, Inches(0.5),
             title, font=HEADER_FONT, size=18, color=INK, bold=True)
    add_text(slide, text_x, top + Inches(0.95), text_w, height - Inches(1.0),
             body, font=BODY_FONT, size=13, color=INK, line_spacing=1.4)

card_w = Inches(3.95)
card_h = Inches(4.4)
card_y = Inches(1.7)
ask_card(s, Inches(0.7),  card_y, card_w, card_h, "1",
         "Mould specifications",
         "Laid density of the mould used for dataset 10, ideally with the "
         "chain-line spacing and a note on which historical mould it was "
         "modelled on.")
ask_card(s, Inches(4.75), card_y, card_w, card_h, "2",
         "View on within-stock variation",
         "Is the 10 – 71 % spread we see across non-calibrated folios "
         "plausibly within-stock variation, or does it suggest something "
         "we should look harder at?")
ask_card(s, Inches(8.80), card_y, card_w, card_h, "3",
         "Acknowledgements",
         "Permission to name you, the Manuscripts Lab, and the Mapping "
         "Paper project in the acknowledgements of the IIB report.")
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
         "Plus anything else you would like the project to be aware of.",
         font=BODY_FONT, size=14, color=INK_LIGHT, align=PP_ALIGN.CENTER)


# ==========================================================================
# Slide 10: thank you
# ==========================================================================
s = add_slide()
add_bg(s, TERRACOTTA)
add_text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.4),
         "Thank you",
         font=HEADER_FONT, size=72, color=SAND, bold=True)
add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6),
         "for the samples, the access, and the manual counts on f5v and f9v.",
         font=HEADER_FONT, size=20, color=SAND)
add_text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.6),
         "Discussion + Q & A",
         font=HEADER_FONT, size=18, color=SAND)
add_rect(s, Inches(0), Inches(0), Inches(0.83), SLIDE_H, SAND)
add_text(s, Inches(0.05), Inches(0.4), Inches(0.7), Inches(0.4),
         "10", font=HEADER_FONT, size=22, color=TERRACOTTA, bold=True)


prs.save(OUT)
print(f"wrote {OUT.relative_to(ROOT)}  ({OUT.stat().st_size/1024/1024:.2f} MB, "
      f"{len(prs.slides)} slides)")
